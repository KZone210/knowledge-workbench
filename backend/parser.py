# -*- coding: utf-8 -*-
"""文档解析模块：支持 PDF / Word / Excel / PPT / Markdown / TXT / HTML / 图片OCR 的文本提取。"""
import os
import re
import html
from html.parser import HTMLParser

SUPPORTED_EXTS = {
    ".pdf", ".docx", ".md", ".markdown", ".txt", ".html", ".htm",
    ".xlsx", ".xls", ".pptx",
    ".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tiff", ".tif", ".gif",
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tiff", ".tif", ".gif"}


class _TextExtractor(HTMLParser):
    """提取 HTML 纯文本（保留段落换行）。"""

    def __init__(self):
        super().__init__()
        self._chunks = []
        self._skip = 0
        self._block_tags = {"p", "div", "br", "li", "h1", "h2", "h3", "h4", "h5", "h6", "tr", "section", "article"}

    def handle_starttag(self, tag, attrs):
        if tag in {"script", "style", "head", "noscript"}:
            self._skip += 1
        if tag in self._block_tags:
            self._chunks.append("\n")

    def handle_endtag(self, tag):
        if tag in {"script", "style", "head", "noscript"}:
            self._skip = max(0, self._skip - 1)
        if tag in self._block_tags:
            self._chunks.append("\n")

    def handle_data(self, data):
        if self._skip == 0 and data.strip():
            self._chunks.append(data)

    def text(self):
        raw = "".join(self._chunks)
        raw = re.sub(r"[ \t\r\f\v]+", " ", raw)
        raw = re.sub(r"\n{3,}", "\n\n", raw)
        return raw.strip()


def _parse_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    pages = []
    for page in reader.pages:
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t.strip():
            pages.append(t)
    return "\n\n".join(pages)


def _parse_docx(path):
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join([c for c in cells if c])
            if line:
                parts.append(line)
    return "\n\n".join(parts)


def _parse_text_file(path):
    for enc in ("utf-8", "gb18030", "utf-16"):
        try:
            with open(path, "r", encoding=enc, errors="strict") as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    # 最终兜底：忽略错误
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _parse_xlsx(path):
    """Excel .xlsx：按工作表提取非空单元格（行内用 | 分隔）。"""
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True, read_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows_text.append(" | ".join(cells))
            if rows_text:
                parts.append(f"[{ws.title}]\n" + "\n".join(rows_text))
    finally:
        wb.close()
    return "\n\n".join(parts)


def _parse_xls(path):
    """Excel .xls（老格式）：xlrd 提取。"""
    import xlrd
    book = xlrd.open_workbook(path)
    parts = []
    for sheet in book.sheets():
        rows_text = []
        for r in range(sheet.nrows):
            cells = [str(sheet.cell_value(r, c)).strip() for c in range(sheet.ncols) if str(sheet.cell_value(r, c)).strip()]
            if cells:
                rows_text.append(" | ".join(cells))
        if rows_text:
            parts.append(f"[{sheet.name}]\n" + "\n".join(rows_text))
    return "\n\n".join(parts)


def _parse_pptx(path):
    """PowerPoint .pptx：逐页提取文本框与表格文本。"""
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        texts.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


_ocr_engine = None


def _get_ocr():
    """RapidOCR 懒加载单例；不可用时返回 None（图片仍可入库，仅无文字）。"""
    global _ocr_engine
    if _ocr_engine is None:
        try:
            from rapidocr_onnxruntime import RapidOCR
            _ocr_engine = RapidOCR()
        except Exception:
            _ocr_engine = False
    return _ocr_engine or None


def _parse_image(path):
    """图片 OCR：识别图中文字（png/jpg/bmp/webp/tiff/gif）。"""
    from PIL import Image
    with Image.open(path) as im:
        im.load()
    engine = _get_ocr()
    if not engine:
        return ""
    try:
        result, _ = engine(str(path))
        if result:
            lines = [item[1] for item in result if item and len(item) > 1 and item[1]]
            return "\n".join(lines)
    except Exception:
        return ""
    return ""


def _parse_html(path):
    raw = _parse_text_file(path)
    parser = _TextExtractor()
    try:
        parser.feed(raw)
        text = parser.text()
    except Exception:
        text = re.sub(r"<[^>]+>", " ", raw)
        text = html.unescape(text)
    return text


def extract_text(path):
    """从文件提取纯文本，返回 (text, title)。"""
    ext = os.path.splitext(path)[1].lower()
    base = os.path.splitext(os.path.basename(path))[0]
    if ext == ".pdf":
        text = _parse_pdf(path)
    elif ext == ".docx":
        text = _parse_docx(path)
    elif ext in (".xlsx",):
        text = _parse_xlsx(path)
    elif ext == ".xls":
        text = _parse_xls(path)
    elif ext == ".pptx":
        text = _parse_pptx(path)
    elif ext in IMAGE_EXTS:
        text = _parse_image(path)
    elif ext in (".md", ".markdown"):
        text = _parse_text_file(path)
        text = re.sub(r"^#{1,6}\s*", "", text, flags=re.M)  # 去掉标题井号
        text = re.sub(r"```.*?```", " ", text, flags=re.S)  # 去掉代码块标记
        text = re.sub(r"!\[.*?\]\(.*?\)", " ", text)        # 去图片语法
        text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)  # 链接转文字
        text = re.sub(r"[*_>`~|#-]", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
    elif ext in (".html", ".htm"):
        text = _parse_html(path)
    else:  # .txt 及其他
        text = _parse_text_file(path)

    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = text.strip()

    # 提取标题：优先取文档内第一个像标题的行
    title = _guess_title(text) or base
    return text, title


def _guess_title(text):
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    if not lines:
        return ""
    first = lines[0]
    # 优先取第一句（按句末标点切分）作为标题
    first_sent = re.split(r"[。！？!?]", first)[0].strip()
    if first_sent and len(first_sent) <= 40:
        return first_sent
    if len(first) <= 60:
        return first
    for line in lines[:8]:
        if len(line) <= 60 and not line.endswith(("。", "！", "？")):
            return line
    return lines[0][:40]
